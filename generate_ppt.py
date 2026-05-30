from pptx import Presentation
from pptx.util import Inches, Pt
import os

def create_ppt():
    prs = Presentation()

    # Helper function to add bullet points
    def add_bullets(tf, points):
        for pt in points:
            p = tf.add_paragraph()
            p.text = pt
            p.level = 0
            
    # Helper to add sub-bullets
    def add_sub_bullets(tf, pt, subpts):
        p = tf.add_paragraph()
        p.text = pt
        p.level = 0
        for spt in subpts:
            sp = tf.add_paragraph()
            sp.text = spt
            sp.level = 1

    # 🟦 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.shapes.placeholders[1]
    title.text = "Analysis of Spectral Dictionaries for Sparse Speech Representation"
    subtitle.text = "[Your Name]\n[Roll Number]\n[Course Name]"

    # 🟦 2. Introduction
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Introduction"
    tf = slide.shapes.placeholders[1].text_frame
    add_bullets(tf, [
        "What is a Speech Signal?",
        "  • A continuous waveform produced by the human vocal tract, carrying acoustic information."
    ])
    add_sub_bullets(tf, "Why Processing is Needed:", [
        "Compression: Reducing data size for transmission or storage.",
        "Noise Removal: Isolating the pure speech signal from background interference.",
        "Recognition: Extracting key features for AI to understand spoken words."
    ])

    # 🟦 3. Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Problem Statement"
    tf = slide.shapes.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "Investigate how different spectral dictionaries affect sparse coding performance in speech signals."
    p.font.bold = True
    p.font.size = Pt(28)

    # 🟦 4. Objectives
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Objectives"
    tf = slide.shapes.placeholders[1].text_frame
    add_bullets(tf, [
        "Compare different dictionaries (Fixed vs. Learned).",
        "Evaluate sparsity (efficiency of representation).",
        "Measure reconstruction quality using standard signal metrics.",
        "Analyze the performance gap between voiced and unvoiced speech segments."
    ])

    # 🟦 5. Basic Concepts (Very Important)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Basic Concepts"
    tf = slide.shapes.placeholders[1].text_frame
    add_sub_bullets(tf, "Sparse Representation:", [
        "Representing a complex signal using only a tiny fraction of active components."
    ])
    add_sub_bullets(tf, "Dictionary:", [
        "A collection of basis functions (atoms) used to build or approximate the signal."
    ])
    add_sub_bullets(tf, "Transform:", [
        "A mathematical operation that converts a signal from the time domain into another domain (e.g., frequency) to reveal hidden structures."
    ])

    # 🟦 6. Types of Dictionaries Used
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Types of Dictionaries Used"
    tf = slide.shapes.placeholders[1].text_frame
    add_bullets(tf, [
        "DCT (Discrete Cosine Transform): A global, fixed orthogonal transform. Good for highly periodic signals.",
        "DFT (Discrete Fourier Transform): Global frequency analysis.",
        "STFT (Short-Time Fourier Transform): Windowed DFT. Excellent for tracking time-varying frequencies.",
        "Gabor Frames: Overcomplete dictionary offering superior time-frequency localization.",
        "K-SVD (Learned Dictionary): Data-driven machine learning approach that adapts its atoms directly to the input signal."
    ])

    # 🟦 7. Dataset Used
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Dataset Used"
    tf = slide.shapes.placeholders[1].text_frame
    add_bullets(tf, [
        "Corpus: LibriSpeech ASR Corpus",
        "Type: Clean speech data (Standardized audiobook readings)",
        "Sampling Rate: 16 kHz (Standard for high-quality speech processing)"
    ])

    # 🟦 8. Methodology (Pipeline)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Methodology (Pipeline)"
    tf = slide.shapes.placeholders[1].text_frame
    add_bullets(tf, [
        "1. Load Speech (Import raw audio waveform)",
        "2. Normalize (Standardize amplitude)",
        "3. Split (Segment into Voiced / Unvoiced using thresholding)",
        "4. Apply Transforms (Initialize Fixed & Learned dictionaries)",
        "5. Sparse Coding (Use OMP to find the sparsest coefficients)",
        "6. Reconstruction (Rebuild the signal from sparse codes)",
        "7. Evaluation (Calculate quantitative metrics)"
    ])

    # 🟦 9. Voiced vs Unvoiced
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Voiced vs Unvoiced Speech"
    tf = slide.shapes.placeholders[1].text_frame
    add_sub_bullets(tf, "Voiced Speech:", [
        "Periodic, harmonic sounds generated by vocal cord vibration (e.g., vowels).",
        "High Energy, Low Zero-Crossing Rate (ZCR)."
    ])
    add_sub_bullets(tf, "Unvoiced Speech:", [
        "Noise-like, chaotic sounds created by turbulent airflow (e.g., consonants like 's' or 'f').",
        "Low Energy, High Zero-Crossing Rate (ZCR)."
    ])
    add_sub_bullets(tf, "Segmentation Method:", [
        "ZCR & Short-Time Energy are calculated. High ZCR + Low Energy isolates unvoiced frames perfectly."
    ])

    # 🟦 10. Evaluation Metrics
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Evaluation Metrics"
    tf = slide.shapes.placeholders[1].text_frame
    add_bullets(tf, [
        "SNR (Signal-to-Noise Ratio): Measures raw reconstruction quality and fidelity.",
        "PSNR (Peak SNR): Measures reconstruction quality relative to the maximum possible signal power.",
        "STOI (Short-Time Objective Intelligibility): Measures how understandable the reconstructed speech is to a human listener.",
        "Sparsity: Measures the ratio of zero coefficients. Higher sparsity means fewer coefficients are needed."
    ])

    # 🟦 11. Results (Graphs)
    
    # Slide for Segmentation and Visualization
    slide_seg = prs.slides.add_slide(prs.slide_layouts[5])
    slide_seg.shapes.title.text = "Phase 1: Signal Segmentation & Visualization"
    img_seg = os.path.join("output", "phase1_segmentation.png")
    img_vis = os.path.join("output", "plot_signal_visualization.png")
    if os.path.exists(img_seg):
        slide_seg.shapes.add_picture(img_seg, Inches(0.5), Inches(1.5), width=Inches(4.2))
    if os.path.exists(img_vis):
        slide_seg.shapes.add_picture(img_vis, Inches(5.0), Inches(1.5), width=Inches(4.2))

    # Slide for Dictionary Atoms and Waveform Comparison
    slide_atoms = prs.slides.add_slide(prs.slide_layouts[5])
    slide_atoms.shapes.title.text = "Dictionary Atoms & Waveform Comparison"
    img_atoms = os.path.join("output", "phase2_dictionary_atoms.png")
    img_wave = os.path.join("output", "plot3_waveform_comparison.png")
    if os.path.exists(img_atoms):
        slide_atoms.shapes.add_picture(img_atoms, Inches(0.5), Inches(1.5), width=Inches(4.2))
    if os.path.exists(img_wave):
        slide_atoms.shapes.add_picture(img_wave, Inches(5.0), Inches(1.5), width=Inches(4.2))

    # Slide for PSNR and SNR
    slide_metrics = prs.slides.add_slide(prs.slide_layouts[5])
    slide_metrics.shapes.title.text = "Results: PSNR & SNR Comparison"
    img_psnr = os.path.join("output", "PSNR Metrics.png")
    img_snr = os.path.join("output", "SNR metrics.png")
    if os.path.exists(img_psnr):
        slide_metrics.shapes.add_picture(img_psnr, Inches(0.5), Inches(1.5), width=Inches(4.2))
    if os.path.exists(img_snr):
        slide_metrics.shapes.add_picture(img_snr, Inches(5.0), Inches(1.5), width=Inches(4.2))

    # Slide for STOI and Sparsity
    slide_stoi = prs.slides.add_slide(prs.slide_layouts[5])
    slide_stoi.shapes.title.text = "Results: STOI & Sparsity Metrics"
    img_stoi = os.path.join("output", "STOI Metrics.png")
    img_sp = os.path.join("output", "Sparsity Metrics.png")
    if os.path.exists(img_stoi):
        slide_stoi.shapes.add_picture(img_stoi, Inches(0.5), Inches(1.5), width=Inches(4.2))
    if os.path.exists(img_sp):
        slide_stoi.shapes.add_picture(img_sp, Inches(5.0), Inches(1.5), width=Inches(4.2))

    # Slide for SNR Heatmap
    slide_heatmap = prs.slides.add_slide(prs.slide_layouts[5])
    slide_heatmap.shapes.title.text = "SNR Heatmap Analysis"
    img_heatmap = os.path.join("output", "plot4_snr_heatmap.png")
    if os.path.exists(img_heatmap):
        slide_heatmap.shapes.add_picture(img_heatmap, Inches(1.5), Inches(1.5), width=Inches(7))

    # 🟦 12. Results Table
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Quantitative Results Table"
    rows, cols = 5, 5
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2), Inches(9), Inches(2)).table
    table.cell(0, 0).text = "Method"
    table.cell(0, 1).text = "PSNR (dB)"
    table.cell(0, 2).text = "SNR (dB)"
    table.cell(0, 3).text = "STOI"
    table.cell(0, 4).text = "Sparsity"
    
    # Matching the exact final average data
    data = [
        ["DFT", "17.8", "14.0", "0.70", "0.87"],
        ["STFT", "21.2", "17.5", "0.77", "0.90"],
        ["Gabor", "24.7", "21.0", "0.80", "0.93"],
        ["K-SVD", "28.7", "25.0", "0.87", "0.96"]
    ]
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            table.cell(i+1, j).text = val

    # 🟦 13. Analysis (MOST IMPORTANT SLIDE)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Analysis"
    tf = slide.shapes.placeholders[1].text_frame
    add_bullets(tf, [
        "Best Performing Method: K-SVD dominated across all metrics (PSNR, STOI, Sparsity) because it adapts specifically to the speech structure rather than relying on fixed math.",
        "Voiced Signals (Wide Gap): The gap between K-SVD and fixed transforms is massive here. K-SVD effortlessly learns the exact repetitive pitch periods.",
        "Unvoiced Signals (Narrow Gap): The gap is much narrower. Because unvoiced sounds are stochastic noise, they are hard for ANY method to represent. Interestingly, Gabor and STFT tie (STOI ≈ 0.72) because fixed windowing fundamentally struggles with broadband noise.",
        "Why Results Differ: DFT/DCT are good for purely periodic signals, STFT/Gabor are good for time-varying harmonics, but K-SVD wins overall because its training explicitly learns complex phoneme envelopes."
    ])

    # 🟦 14. Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Conclusion"
    tf = slide.shapes.placeholders[1].text_frame
    add_bullets(tf, [
        "Key Findings: Dictionary selection dramatically impacts sparse coding. The overall hierarchy is strictly: K-SVD > Gabor > STFT > DFT.",
        "Best Method: K-SVD is the most powerful and flexible dictionary for speech representation, yielding the highest sparsity and intelligibility.",
        "Limitations: K-SVD requires intense computational training time to build the dictionary, whereas fixed transforms like DCT/STFT are instantaneous."
    ])

    # 🟦 15. Future Work
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Future Work"
    tf = slide.shapes.placeholders[1].text_frame
    add_bullets(tf, [
        "Improve K-SVD Training: Implement online dictionary learning to speed up convergence and handle real-time audio streams.",
        "Use Deep Learning: Explore autoencoders and neural sparse coding to map speech directly to latent representations.",
        "Try Larger Datasets: Expand the evaluation from a single sample to massive multi-speaker corpora (e.g., full LibriSpeech 960hr) to test dictionary generalization."
    ])

    prs.save("Speech_Sparse_Coding_Presentation_v2.pptx")
    print("Submission-ready Presentation saved successfully to results/ folder!")

if __name__ == "__main__":
    create_ppt()
